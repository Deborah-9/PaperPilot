import logging
from typing import Dict, List, Optional
import os
from datetime import datetime
import tempfile
from pathlib import Path
import fitz  # PyMuPDF for PDF handling
import docx  # python-docx for DOCX handling
from pptx import Presentation  # python-pptx for PowerPoint
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup, ParseMode, ChatAction, CallbackQuery
from telegram.ext import CallbackContext
import google.generativeai as genai
from io import BytesIO
import re
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Enable logging
logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)
logger = logging.getLogger(__name__)

# Configure Gemini AI
GOOGLE_API_KEY = os.getenv('GOOGLE_API_KEY')
if not GOOGLE_API_KEY:
    raise ValueError("❌ Google API Key not found in environment variables!")

# Initialize Gemini
genai.configure(api_key=GOOGLE_API_KEY)

# Fun messages for non-academic documents
NON_ACADEMIC_MESSAGES = [
    "🎭 Whoopsie! Looks like this document is playing hide and seek with academic content! Try sending something more scholarly! 📚",
    "🌟 Plot twist: This document seems more casual than a Friday night pizza! Let's keep it academic! 🎓",
    "🎪 While I appreciate creativity, this looks more like weekend reading than research material! 📖",
    "🚀 Houston, we have a situation! This document seems to have wandered away from Academia Street! 🗺️",
    "🎨 Beautiful document, but I'm more into the academic masterpieces! Care to share some research instead? 🔬",
    "🎮 This looks fun, but I'm more of a scholarly papers kind of bot! Got any of those? 📝",
]

# Keywords indicating academic content
ACADEMIC_INDICATORS = {
    'strong': [
        r'abstract', r'introduction', r'methodology', r'conclusion',
        r'references', r'citation[s]?', r'bibliography', r'hypothesis',
        r'research', r'analysis', r'study', r'experiment[s]?',
        r'data', r'results', r'discussion', r'findings'
    ],
    'moderate': [
        r'figure[s]?', r'table[s]?', r'equation[s]?', r'theory',
        r'model[s]?', r'algorithm[s]?', r'method[s]?', r'framework',
        r'approach', r'implementation', r'evaluation'
    ]
}

class DocumentHandler:
    """Handles document processing and analysis using Gemini AI."""

    def __init__(self):
        """Initialize DocumentHandler with its own Gemini model instance."""
        try:
            # Initialize Gemini model with advanced configuration
            generation_config = {
                'temperature': 0.7,
                'top_p': 0.9,
                'top_k': 40,
                'max_output_tokens': 2048,
            }

            safety_settings = [
                {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
                {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
                {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
                {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"}
            ]

            self.model = genai.GenerativeModel(
                model_name='gemini-1.5-pro',
                generation_config=generation_config,
                safety_settings=safety_settings
            )

            logger.info("🤖 Gemini AI model initialized successfully!")

        except Exception as e:
            logger.error(f"❌ Failed to initialize Gemini AI model: {str(e)}")
            raise

        self.user_documents: Dict[int, Dict] = {}
        self.max_chunk_size = 30000

        self.system_prompt = """
        You are an academic research assistant specialized in analyzing academic documents.
        Your responses should be:
        1. Scholarly and precise
        2. Well-structured with clear sections
        3. Enhanced with relevant emojis for better readability
        4. Focused on academic content and research implications
        """

        logger.info("📚 DocumentHandler initialized and ready!")

    def start(self, update: Update, context: CallbackContext) -> None:
        """Send welcome message when document handler is first used."""
        update.message.reply_text(
            "🌟 *Welcome to PaperPilot's Document Analysis!*\n\n"
            "I can help you analyze academic documents including:\n"
            "📄 PDF files\n"
            "📝 DOCX files\n"
            "📊 PPTX presentations\n"
            "📰 TXT files\n\n"
            "Just send me an academic document and I'll:\n"
            "✨ Generate summaries\n"
            "🔍 Provide detailed analysis\n"
            "❓ Answer your questions\n"
            "📊 Extract key points\n"
            "🔗 Find related papers\n"
            "🎯 Identify research gaps\n\n"
            "_Send me an academic document to get started!_ 🚀",
            parse_mode=ParseMode.MARKDOWN
        )

    def handle_document(self, update: Update, context: CallbackContext) -> None:
        """Main handler for document messages."""
        chat_id = update.effective_chat.id
        message = update.message

        # Send "typing" action
        context.bot.send_chat_action(chat_id=chat_id, action=ChatAction.TYPING)

        # Get document file
        doc_file = context.bot.get_file(message.document.file_id)
        file_extension = Path(message.document.file_name).suffix.lower()

        # Check supported formats
        if file_extension not in ['.pdf', '.txt', '.docx', '.pptx']:
            message.reply_text(
                "🎭 Oops! I can only read PDF, TXT, DOCX, and PPTX files!\n"
                "Send me one of those, and I'll work my magic! ✨",
                parse_mode=ParseMode.MARKDOWN
            )
            return

        # Create temp file
        with tempfile.NamedTemporaryFile(suffix=file_extension, delete=False) as temp_file:
            # Download document
            doc_file.download(custom_path=temp_file.name)

            # Extract text based on file type
            text_content = self._extract_text(temp_file.name, file_extension)

            # Check if document is academic
            if not self._is_academic_document(text_content):
                message.reply_text(
                    self._get_random_non_academic_message(),
                    parse_mode=ParseMode.MARKDOWN
                )
                os.unlink(temp_file.name)
                return

            # Process and store document
            doc_info = self._process_document(text_content, message.document.file_name)
            self.user_documents[chat_id] = doc_info

            # Create interactive keyboard
            keyboard = [
                [InlineKeyboardButton("📝 Summary", callback_data="doc_summary"),
                 InlineKeyboardButton("🔍 Detailed Analysis", callback_data="doc_analysis")],
                [InlineKeyboardButton("❓ Ask Question", callback_data="doc_question"),
                 InlineKeyboardButton("📊 Key Points", callback_data="doc_keypoints")],
                [InlineKeyboardButton("📚 Related Papers", callback_data="doc_related"),
                 InlineKeyboardButton("🎯 Research Gap", callback_data="doc_gaps")]
            ]
            reply_markup = InlineKeyboardMarkup(keyboard)

            # Send welcome message
            message.reply_text(
                f"🌟 *Document Analysis Ready!*\n\n"
                f"I've processed your document: `{message.document.file_name}`\n\n"
                f"What would you like to know about it? Choose an option below or simply ask me anything! 🤓",
                parse_mode=ParseMode.MARKDOWN,
                reply_markup=reply_markup
            )

    def _extract_text(self, file_path: str, extension: str) -> str:
        """Extract text from different document formats."""
        try:
            if extension == '.pdf':
                return self._extract_pdf_text(file_path)
            elif extension == '.txt':
                return self._extract_txt_text(file_path)
            elif extension == '.docx':
                return self._extract_docx_text(file_path)
            elif extension == '.pptx':
                return self._extract_pptx_text(file_path)
            else:
                raise ValueError(f"Unsupported file format: {extension}")
        except Exception as e:
            logger.error(f"Error extracting text from {extension} file: {str(e)}")
            raise

    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF files."""
        text = []
        with fitz.open(file_path) as doc:
            for page in doc:
                text.append(page.get_text())
        return "\n".join(text)

    def _extract_txt_text(self, file_path: str) -> str:
        """Extract text from TXT files."""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()

    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from DOCX files."""
        doc = docx.Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PPTX files."""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def _is_academic_document(self, text: str) -> bool:
        """Check if document appears to be academic."""
        text_lower = text.lower()
        strong_matches = sum(1 for pattern in ACADEMIC_INDICATORS['strong']
                           if re.search(pattern, text_lower))
        moderate_matches = sum(1 for pattern in ACADEMIC_INDICATORS['moderate']
                             if re.search(pattern, text_lower))
        return (strong_matches >= 3) or (strong_matches >= 2 and moderate_matches >= 3)

    def _get_random_non_academic_message(self) -> str:
        """Return random fun message for non-academic docs."""
        return NON_ACADEMIC_MESSAGES[int(datetime.now().timestamp()) % len(NON_ACADEMIC_MESSAGES)]

    def _process_document(self, text: str, filename: str) -> Dict:
        """Process and analyze document content."""
        chunks = self._split_into_chunks(text)
        prompt = f"""
        Analyze this academic document and provide:
        1. A brief summary
        2. Main topics covered
        3. Key findings or arguments
        4. Methodology used (if applicable)
        5. Potential research gaps or future work

        Document text chunk (1/{len(chunks)}):
        {chunks[0]}
        """

        response = self.model.generate_content(prompt)

        return {
            'filename': filename,
            'chunks': chunks,
            'initial_analysis': response.text,
            'timestamp': datetime.now().isoformat(),
            'chat_history': []
        }

    def _split_into_chunks(self, text: str) -> List[str]:
        """Split text into manageable chunks for processing."""
        chunks = []
        current_chunk = []
        current_length = 0

        for paragraph in text.split('\n'):
            if current_length + len(paragraph) > self.max_chunk_size:
                if current_chunk:
                    chunks.append('\n'.join(current_chunk))
                current_chunk = [paragraph]
                current_length = len(paragraph)
            else:
                current_chunk.append(paragraph)
                current_length += len(paragraph)

        if current_chunk:
            chunks.append('\n'.join(current_chunk))

        return chunks

    def handle_document_query(self, update: Update, context: CallbackContext) -> None:
        """Handle user queries about the document."""
        query = update.callback_query
        chat_id = update.effective_chat.id

        if chat_id not in self.user_documents:
            query.message.reply_text(
                "🎭 Oops! I don't see any active document to analyze.\n"
                "Send me an academic document first! 📚",
                parse_mode=ParseMode.MARKDOWN
            )
            return

        doc_info = self.user_documents[chat_id]
        action = query.data

        if action == "doc_summary":
            self._generate_summary(query, doc_info)
        elif action == "doc_analysis":
            self._generate_detailed_analysis(query, doc_info)
        elif action == "doc_question":
            self._handle_question_mode(query, doc_info)
        elif action == "doc_keypoints":
            self._generate_key_points(query, doc_info)
        elif action == "doc_related":
            self._find_related_papers(query, doc_info)
        elif action == "doc_gaps":
            self._analyze_research_gaps(query, doc_info)
        elif action == "doc_main_menu":
            self._show_main_menu(query, doc_info)

    def handle_text_query(self, update: Update, context: CallbackContext) -> None:
        """Handle text questions about the active document."""
        chat_id = update.effective_chat.id
        message = update.message

        if chat_id not in self.user_documents:
            message.reply_text(
                "🎭 I'd love to help, but I don't see any active document!\n"
                "Send me an academic document first, and I'll answer all your questions! 📚",
                parse_mode=ParseMode.MARKDOWN
            )
            return

        doc_info = self.user_documents[chat_id]
        doc_info['chat_history'].append({
            'role': 'user',
            'content': message.text,
            'timestamp': datetime.now().isoformat()
        })

        prompt = self._create_chat_prompt(doc_info, message.text)
        response = self.model.generate_content(prompt)

        doc_info['chat_history'].append({
            'role': 'assistant',
            'content': response.text,
            'timestamp': datetime.now().isoformat()
        })

        message.reply_text(response.text, parse_mode=ParseMode.MARKDOWN)

    def _create_chat_prompt(self, doc_info: Dict, question: str) -> str:
        """Create context-aware prompt for chat."""
        return f"""
        Based on the academic document "{doc_info['filename']}", please answer:

        {question}

        Previous conversation:
        {self._format_chat_history(doc_info['chat_history'][-5:])}

        Document analysis:
        {doc_info['initial_analysis']}
        """

    def _format_chat_history(self, history: List[Dict]) -> str:
        """Format chat history for context."""
        return "\n".join(f"{msg['role'].upper()}: {msg['content']}" for msg in history)

    def _generate_summary(self, query: CallbackQuery, doc_info: Dict) -> None:
        """Generate and send document summary."""
        query.answer("🎯 Generating summary...")
        prompt = f"""
        Create a concise summary of this academic document.
        Focus on main contributions, methodology, and key findings.
        Format with clear sections and emoji indicators.

        Document: {doc_info['initial_analysis']}
        """

        response = self.model.generate_content(prompt)
        query.message.reply_text(
            f"📚 *Document Summary*\n\n{response.text}",
            parse_mode=ParseMode.MARKDOWN
        )

    def _generate_detailed_analysis(self, query: CallbackQuery, doc_info: Dict) -> None:
        """Generate detailed analysis options."""
        query.answer("🔬 Analyzing in detail...")

        keyboard = [
            [InlineKeyboardButton("📊 Methods", callback_data="analysis_methods"),
             InlineKeyboardButton("🎯 Results", callback_data="analysis_results")],
            [InlineKeyboardButton("💡 Innovation", callback_data="analysis_innovation"),
             InlineKeyboardButton("📈 Impact", callback_data="analysis_impact")],
            [InlineKeyboardButton("🔄 Back", callback_data="doc_main_menu")]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)

        query.message.reply_text(
            "🔍 *Choose an aspect to analyze:*\n\n"
            "• *Methods*: Research methodology\n"
            "• *Results*: Key findings\n"
            "• *Innovation*: Novel contributions\n"
            "• *Impact*: Significance",
            parse_mode=ParseMode.MARKDOWN,
            reply_markup=reply_markup
        )

    def _handle_question_mode(self, query: CallbackQuery, doc_info: Dict) -> None:
        """Enable question mode for the document."""
        query.answer("❓ Question mode activated!")

        keyboard = [
            [InlineKeyboardButton("📝 Example Questions", callback_data="doc_example_questions"),
             InlineKeyboardButton("🎯 Exit Q&A", callback_data="doc_main_menu")]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)

        query.message.reply_text(
            "🤓 *Ask Me Anything Mode Activated!*\n\n"
            "Just type your question about the document, and I'll answer!\n\n"
            "Some things you can ask:\n"
            "• What are the main findings?\n"
            "• Explain the methodology used\n"
            "• What datasets were used?\n"
            "• What are the limitations?\n"
            "• Compare this with existing work\n\n"
            "🎯 _Try asking something specific!_",
            parse_mode=ParseMode.MARKDOWN,
            reply_markup=reply_markup
        )

    def _generate_key_points(self, query: CallbackQuery, doc_info: Dict) -> None:
        """Extract and present key points from the document."""
        query.answer("📊 Extracting key points...")

        prompt = f"""
        Extract the most important key points from this academic document.
        Include:
        1. Major contributions
        2. Critical findings
        3. Novel insights
        4. Important conclusions

        Format them as a bulleted list with emojis.

        Document analysis: {doc_info['initial_analysis']}
        """

        response = self.model.generate_content(prompt)

        keyboard = [
            [InlineKeyboardButton("📈 Visualize Points", callback_data="keypoints_visualize"),
             InlineKeyboardButton("🔍 Expand Point", callback_data="keypoints_expand")],
            [InlineKeyboardButton("🔄 Back", callback_data="doc_main_menu")]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)

        query.message.reply_text(
            f"📊 *Key Points Analysis*\n\n{response.text}",
            parse_mode=ParseMode.MARKDOWN,
            reply_markup=reply_markup
        )

    def _find_related_papers(self, query: CallbackQuery, doc_info: Dict) -> None:
        """Find and suggest related papers."""
        query.answer("🔍 Searching related papers...")

        prompt = f"""
        Extract 3-5 key search terms from this document that would be
        most effective for finding related academic papers. Focus on:
        1. Core concepts
        2. Methodologies
        3. Research areas

        Document: {doc_info['initial_analysis']}
        """

        search_terms = self.model.generate_content(prompt)

        keyboard = []
        for term in search_terms.text.split('\n'):
            if term.strip():
                keyboard.append([
                    InlineKeyboardButton(f"🔍 {term.strip()}",
                                      callback_data=f"search_{term.strip()}")
                ])

        keyboard.append([InlineKeyboardButton("🔄 Back", callback_data="doc_main_menu")])
        reply_markup = InlineKeyboardMarkup(keyboard)

        query.message.reply_text(
            "📚 *Related Papers Search*\n\n"
            "I've extracted key topics from your document.\n"
            "Click on any term to find related papers!\n\n"
            "*Search Terms:*\n" + search_terms.text,
            parse_mode=ParseMode.MARKDOWN,
            reply_markup=reply_markup
        )

    def _analyze_research_gaps(self, query: CallbackQuery, doc_info: Dict) -> None:
        """Analyze and present research gaps and future work."""
        query.answer("🎯 Analyzing research gaps...")

        prompt = f"""
        Analyze this academic document to identify:
        1. Current research gaps
        2. Limitations of the work
        3. Potential future research directions
        4. Unexplored aspects of the topic

        Format the response with clear sections and recommendations.

        Document: {doc_info['initial_analysis']}
        """

        response = self.model.generate_content(prompt)

        keyboard = [
            [InlineKeyboardButton("💡 Suggest Extensions", callback_data="gaps_suggest_extensions"),
             InlineKeyboardButton("🔍 Explore Gap", callback_data="gaps_explore")],
            [InlineKeyboardButton("🔄 Back", callback_data="doc_main_menu")]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)

        query.message.reply_text(
            f"🎯 *Research Gaps Analysis*\n\n{response.text}",
            parse_mode=ParseMode.MARKDOWN,
            reply_markup=reply_markup
        )

    def _show_main_menu(self, query: CallbackQuery, doc_info: Dict) -> None:
        """Show the main document analysis menu."""
        query.answer()
        keyboard = [
            [InlineKeyboardButton("📝 Summary", callback_data="doc_summary"),
             InlineKeyboardButton("🔍 Detailed Analysis", callback_data="doc_analysis")],
            [InlineKeyboardButton("❓ Ask Question", callback_data="doc_question"),
             InlineKeyboardButton("📊 Key Points", callback_data="doc_keypoints")],
            [InlineKeyboardButton("📚 Related Papers", callback_data="doc_related"),
             InlineKeyboardButton("🎯 Research Gap", callback_data="doc_gaps")]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)

        query.edit_message_text(
            f"🌟 *Document Analysis Menu*\n\n"
            f"Current document: `{doc_info['filename']}`\n\n"
            f"What would you like to know about it?",
            parse_mode=ParseMode.MARKDOWN,
            reply_markup=reply_markup
        )

    def handle_analysis_query(self, update: Update, context: CallbackContext) -> None:
        """Handle detailed analysis button callbacks."""
        query = update.callback_query
        chat_id = update.effective_chat.id

        if chat_id not in self.user_documents:
            query.answer("No active document! Send me one first! 📚")
            return

        doc_info = self.user_documents[chat_id]
        action = query.data

        analysis_prompts = {
            'analysis_methods': (
                "🔬 *Methodology Analysis*\n\n"
                "Analyze and explain the research methodology, including:\n"
                "1. Research approach\n"
                "2. Data collection methods\n"
                "3. Analysis techniques\n"
                "4. Validation methods"
            ),
            'analysis_results': (
                "📊 *Results Analysis*\n\n"
                "Analyze the key findings, including:\n"
                "1. Main results\n"
                "2. Statistical significance\n"
                "3. Practical implications\n"
                "4. Comparative analysis"
            ),
            'analysis_innovation': (
                "💡 *Innovation Analysis*\n\n"
                "Identify and analyze innovative aspects:\n"
                "1. Novel contributions\n"
                "2. Technical advancements\n"
                "3. Theoretical contributions\n"
                "4. Methodological innovations"
            ),
            'analysis_impact': (
                "📈 *Impact Analysis*\n\n"
                "Evaluate the research impact:\n"
                "1. Academic significance\n"
                "2. Practical applications\n"
                "3. Industry relevance\n"
                "4. Future implications"
            )
        }

        if action in analysis_prompts:
            query.answer(f"Analyzing {action.split('_')[1]}...")
            prompt = f"{analysis_prompts[action]}\n\nDocument: {doc_info['initial_analysis']}"
            response = self.model.generate_content(prompt)

            keyboard = [[InlineKeyboardButton("🔄 Back to Analysis", callback_data="doc_analysis")]]
            reply_markup = InlineKeyboardMarkup(keyboard)

            query.edit_message_text(
                response.text,
                parse_mode=ParseMode.MARKDOWN,
                reply_markup=reply_markup
            )

        elif action == "doc_main_menu":
            self._show_main_menu(query, doc_info)